import streamlit as st
import math
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import io

# --- KONFIGURASI BARU (FONT 50pt) ---
# Karena font lebih kecil, kita bisa memuat lebih banyak kata.
TARGET_IDEAL = 12   # Sebelumnya 10
BATAS_MAKSIMAL = 18 # Sebelumnya 15

# --- FUNGSI LOGIKA ---
def bagi_secara_adil(teks, batas_maksimal):
    # Menggunakan split(None) agar aman terhadap spasi ganda/tab/nbsp
    kata_kata = teks.split()
    total_kata = len(kata_kata)
    if total_kata <= batas_maksimal: return [teks]
    jumlah_chunk = math.ceil(total_kata / batas_maksimal)
    kata_per_chunk = math.ceil(total_kata / jumlah_chunk)
    potongan = []
    for i in range(0, total_kata, kata_per_chunk):
        chunk = kata_kata[i : i + kata_per_chunk]
        potongan.append(" ".join(chunk))
    return potongan

def proses_kalimat_final(teks_lengkap):
    # Split berdasarkan / tapi lebih toleran terhadap spasi
    raw_phrases = [p.strip() for p in re.split(r'\s*/\s*', teks_lengkap) if p.strip()]
    final_slides = []
    current_slide_phrases = []
    current_word_count = 0
    
    for phrase in raw_phrases:
        phrase_word_count = len(phrase.split())
        
        if phrase_word_count > BATAS_MAKSIMAL:
            if current_slide_phrases:
                final_slides.append(" / ".join(current_slide_phrases) + " /")
                current_slide_phrases = []
                current_word_count = 0
            sub_chunks = bagi_secara_adil(phrase, TARGET_IDEAL)
            for i, chunk in enumerate(sub_chunks):
                if i == len(sub_chunks) - 1:
                    current_slide_phrases.append(chunk)
                    current_word_count = len(chunk.split())
                else:
                    final_slides.append(chunk + " /")
            continue

        potensi_total = current_word_count + phrase_word_count
        masuk = False
        if potensi_total <= TARGET_IDEAL:
            masuk = True
        elif potensi_total <= BATAS_MAKSIMAL:
            if "//" in phrase or ")" in phrase: masuk = True
            else: masuk = False
        else: masuk = False
            
        if masuk:
            current_slide_phrases.append(phrase)
            current_word_count += phrase_word_count
        else:
            if current_slide_phrases:
                final_slides.append(" / ".join(current_slide_phrases) + " /")
            current_slide_phrases = [phrase]
            current_word_count = phrase_word_count
            
    if current_slide_phrases:
        final_slides.append(" / ".join(current_slide_phrases))
    return final_slides

# --- FUNGSI PEMBUAT PPT ---
def generate_pptx_binary(naskah_text):
    prs = Presentation()
    # Format 4:3 (10 inci x 7.5 inci)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    tanda_sumber = chr(91) + "source:" 
    tanda_tutup = chr(93)

    lines = naskah_text.split('\n')
    for line in lines:
        text_raw = line.strip()
        
        # --- DEEP CLEANING ---
        text_raw = text_raw.replace('\u200b', '')
        text_raw = text_raw.replace('\xa0', ' ')
        
        if tanda_sumber in text_raw:
            parts = text_raw.split(tanda_tutup)
            text_clean = parts[1].strip() if len(parts) > 1 else text_raw
        else: text_clean = text_raw
        
        if not text_clean: continue

        # --- AUTO PUNCTUATION ---
        text_clean = text_clean.replace(".", " //")
        text_clean = text_clean.replace(",", " /")
        
        has_end_marker = "//" in text_clean
        text_processing = text_clean.replace("//", " //")

        # --- SMART SPLIT ---
        pecahan = text_processing.split(None, 1)
        
        nama_pembicara = ""
        isi_pesan = text_processing

        if len(pecahan) > 1:
            kata_pertama = pecahan[0]
            # Syarat Nama: Huruf Besar Semua & > 1 Huruf
            if kata_pertama.isupper() and len(kata_pertama) > 1:
                if not kata_pertama.startswith("/"):
                    nama_pembicara = kata_pertama
                    isi_pesan = pecahan[1]

        list_slide_text = proses_kalimat_final(isi_pesan)

        for i, text_slide in enumerate(list_slide_text):
            is_last_slide = (i == len(list_slide_text) - 1)
            text_slide = text_slide.strip()
            
            if is_last_slide:
                if text_slide.endswith("/"): text_slide = text_slide[:-1].strip()
                if has_end_marker and not text_slide.endswith("//"): text_slide += " //"
            else:
                if not text_slide.endswith("/"): text_slide += " /"

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            
            # Area Text Box
            left = Inches(0.2); top = Inches(0.2); width = Inches(9.6); height = Inches(7.1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE 
            tf.word_wrap = True 
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER 

            # --- SETTING FONT 50pt ---
            if nama_pembicara:
                run_nama = p.add_run()
                run_nama.text = nama_pembicara + "\n"
                run_nama.font.name = "Arial Black"
                run_nama.font.size = Pt(50) # UBAH KE 50
                run_nama.font.bold = True
                run_nama.font.color.rgb = RGBColor(255, 255, 0)

            sisa_text = text_slide
            while True:
                idx_buka = sisa_text.find("(")
                if idx_buka == -1:
                    if sisa_text:
                        run = p.add_run()
                        run.text = sisa_text
                        run.font.name = "Arial Black"
                        run.font.size = Pt(50) # UBAH KE 50
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                    break
                idx_tutup = sisa_text.find(")", idx_buka)
                if idx_tutup == -1: 
                    run = p.add_run()
                    run.text = sisa_text
                    run.font.name = "Arial Black"
                    run.font.size = Pt(50) # UBAH KE 50
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    break
                
                text_sebelum = sisa_text[:idx_buka]
                if text_sebelum:
                    run = p.add_run()
                    run.text = text_sebelum
                    run.font.name = "Arial Black"
                    run.font.size = Pt(50) # UBAH KE 50
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                
                text_aksi = sisa_text[idx_buka : idx_tutup+1]
                run_merah = p.add_run()
                run_merah.text = text_aksi
                run_merah.font.name = "Arial Black"
                run_merah.font.size = Pt(50) # UBAH KE 50
                run_merah.font.bold = True
                run_merah.font.color.rgb = RGBColor(255, 0, 0)
                sisa_text = sisa_text[idx_tutup+1:]

    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# --- TAMPILAN WEBSITE ---
st.set_page_config(page_title="Prompter Maker 50pt", page_icon="ðŸ“º")
st.title("ðŸ“º TV Prompter Generator (Font 50pt)")
st.write("Fitur: Font 50pt (Lebih Padat), Host Kuning, Aksi Merah, Anti-Bug Spasi.")

naskah_input = st.text_area("Masukkan Naskah:", height=300, placeholder="HOST Halo pemirsa. (SENYUM) Apa kabar?")

if naskah_input:
    if st.button("Buat File PowerPoint"):
        with st.spinner("Sedang memproses..."):
            file_ppt = generate_pptx_binary(naskah_input)
            st.success("Selesai! Font ukuran 50pt diterapkan.")
            st.download_button(
                label="ðŸ“¥ Download PPT Prompter",
                data=file_ppt,
                file_name="Prompter_50pt.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )